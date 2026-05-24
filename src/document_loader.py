from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger("noguake.doc")

SUPPORTED = {".pdf", ".pptx", ".ppt", ".docx", ".txt", ".md"}


def load_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _load_pdf(path)
    if suffix in (".pptx", ".ppt"):
        return _load_pptx(path)
    if suffix == ".docx":
        return _load_docx(path)
    if suffix in (".txt", ".md"):
        return path.read_text(encoding="utf-8", errors="ignore")
    raise ValueError(f"不支持的文件类型: {path}")


def _load_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return "\n".join(parts)


def _load_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            parts.append(f"[幻灯片 {i}]\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def _load_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def collect_documents(data_dir: Path) -> list[tuple[Path, str]]:
    """返回 (文件路径, 全文) 列表。"""
    results: list[tuple[Path, str]] = []
    if not data_dir.exists():
        return results
    for path in sorted(data_dir.rglob("*")):
        if not path.is_file() or path.suffix.lower() not in SUPPORTED:
            continue
        try:
            text = load_text(path).strip()
            if len(text) >= 30:
                results.append((path, text))
        except Exception as e:
            logger.warning("跳过文件 %s: %s", path.name, e)
    return results
